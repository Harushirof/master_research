# -*- coding: utf-8 -*-
"""
slidekit — 齋藤研の卒論/中間発表フォーマットでスライドを作るヘルパー。

考え方（重要）:
  デザイン（濃紺アクセント・見出し・フォント・■箇条書き）は assets/template.pptx の
  スレイドマスター/レイアウト '1行タイトル+本文' に入っている。手作りで再現しようとすると
  必ずズレるので、**テンプレートを開き、そのレイアウトでスライドを足してプレースホルダを
  埋める** ことで本物の体裁をそのまま継承する。

使い方:
  from slidekit import Deck
  d = Deck()                                  # テンプレートを開く
  d.set_title(["水晶振動発振器の位相同期制御", "による高純度信号作成"],
              student="08-XXXXXX 福井 晴士郎", date="2026/06/24")
  d.set_index(["01 研究背景","02 本研究の目標","03 先行研究","04 研究方法",
               "05 実験と結果","06 考察と今後の展望","Appendix"])
  d.content("01 研究背景：…", ["■の本文1","■の本文2",(1,"●下位の本文")],
            images=["/abs/fig.png"], caps=["図の説明"], layout="right")
  d.content("04-5 制御ロジック", [...], images=[a,b], layout="below2")
  d.save("out.pptx")                          # 仕上げ（卒論の元コンテンツを削除して保存）

注意:
  - bullets は文字列（第1階層=■）か (level, 文字列) のタプル。level1=●, level2=→ は
    レイアウトの箇条書きスタイルに従う。
  - images は絶対パス推奨。layout は 'right'(右に1枚) / 'below2'(下に2枚) / 'below'(下に1枚大)。
  - 図がまだ無い箇所は placeholder=... で灰色の枠（「○○：作成予定」）を置ける。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.normpath(os.path.join(HERE, "..", "assets", "template.pptx"))
CONTENT_LAYOUT = "1行タイトル+本文"     # 本文スライドのレイアウト名（デザインの本体）
NAVY = RGBColor(0x00, 0x0F, 0x78)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
JP = "Yu Gothic"
SW = 10.833  # スライド幅[in]（卒論テンプレ＝9906000EMU）


def _ea(run, font=JP):
    """日本語フォントを latin/ea 両方に設定（python-pptx に get_or_add_ea が無いため手動）。"""
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    latin = rPr.get_or_add_latin(); latin.set("typeface", font)
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {}); latin.addnext(ea)
    ea.set("typeface", font)


def _fit(path, maxw, maxh):
    iw, ih = Image.open(path).size; ar = iw / ih
    w = maxw; h = w / ar
    if h > maxh:
        h = maxh; w = h * ar
    return w, h


class Deck:
    def __init__(self, template=None, keep_first=2):
        self.prs = Presentation(template or TEMPLATE)
        self._n0 = len(self.prs.slides._sldIdLst)   # 開いた時点のスライド数
        self._keep = keep_first   # 先頭から残す枚数（表紙・Index・参照スライド等）。残り(=元の本文)は save() で削除
        self._page = 1

    # ---------- 低レベル ----------
    def _layout(self, name=CONTENT_LAYOUT):
        for l in self.prs.slide_layouts:
            if l.name == name:
                return l
        return self.prs.slide_layouts[0]

    def _find_tb(self, slide, substr):
        for sh in slide.shapes:
            if sh.has_text_frame and substr in sh.text_frame.text:
                return sh
        return None

    def _set_tf(self, tf, lines):
        """テキストフレームの文字を入れ替え（元の1行目の書式＝サイズ/色/太字を継承）。"""
        size = bold = color = None
        for p in tf.paragraphs:
            for r in p.runs:
                size = r.font.size; bold = r.font.bold
                try:
                    if r.font.color and r.font.color.type is not None:
                        color = r.font.color.rgb
                except Exception:
                    pass
                break
            if size is not None:
                break
        tf.clear()
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = ln
            if size is not None:
                r.font.size = size
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color
            _ea(r)

    def _pic(self, slide, path, l, t, maxw, maxh):
        if not (path and os.path.isfile(path)):
            self._ph(slide, l, t, maxw, maxh, "[図ファイル無し: %s]" % os.path.basename(str(path)))
            return
        w, h = _fit(path, maxw, maxh)
        slide.shapes.add_picture(path, Inches(l + (maxw - w) / 2), Inches(t + (maxh - h) / 2), Inches(w), Inches(h))

    def _ph(self, slide, l, t, w, h, text):
        sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
        sp.line.color.rgb = GRAY; sp.line.width = Pt(1); sp.shadow.inherit = False
        tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text; r.font.size = Pt(12); r.font.color.rgb = GRAY; _ea(r)

    def _cap(self, slide, l, t, w, text):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.3))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text; r.font.size = Pt(9); r.font.color.rgb = GRAY; _ea(r)

    # ---------- 表紙・Index ----------
    def set_title(self, title_lines, student=None, date=None, lab=None):
        """表紙（スライド0）の文字を差し替え。指定したものだけ変更。"""
        s = self.prs.slides[0]
        tb = self._find_tb(s, "位相ノイズ測定")            # 元タイトル
        if tb and title_lines:
            self._set_tf(tb.text_frame, title_lines if isinstance(title_lines, list) else [title_lines])
        if student:
            tb = self._find_tb(s, "福井") or self._find_tb(s, "08-")
            if tb:
                self._set_tf(tb.text_frame, [student])
        if date:
            tb = self._find_tb(s, "2025/02/05") or self._find_tb(s, "2025")
            if tb:
                self._set_tf(tb.text_frame, [date])
        if lab:
            tb = self._find_tb(s, "研究室")
            if tb:
                self._set_tf(tb.text_frame, [lab])

    def set_index(self, items):
        """Index（スライド1）のシェブロンを items で作り直す（濃紺・白文字）。"""
        s = self.prs.slides[1]
        for sh in list(s.shapes):
            if sh.shape_type == 6:   # GROUP（既存シェブロン）
                sh._element.getparent().remove(sh._element)
        n = len(items); h = 0.62; gap = 0.16
        total = n * h + (n - 1) * gap
        top = max(1.55, (7.5 - total) / 2 + 0.2)
        for i, it in enumerate(items):
            ch = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(1.3), Inches(top + i * (h + gap)), Inches(7.4), Inches(h))
            ch.fill.solid(); ch.fill.fore_color.rgb = NAVY; ch.line.fill.background(); ch.shadow.inherit = False
            tf = ch.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = False
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = "　" + it; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; _ea(r)

    # ---------- 本文スライド ----------
    def content(self, title, bullets, images=None, caps=None, layout="right", placeholder=None):
        s = self.prs.slides.add_slide(self._layout())
        s.shapes.title.text = title                          # 見出し（濃紺アクセントはレイアウト由来）
        body = s.placeholders[1]
        # 本文枠の位置・幅・高さを明示する。継承のまま .height だけ設定すると幅が欠落し
        # （cx が無い xfrm になり）本文が描画されない不具合があるため、4値とも必ず与える。
        right_col = (placeholder is not None) or (layout == "right" and images)
        body.left = Inches(0.30); body.top = Inches(1.00)
        if right_col:
            body.width = Inches(4.95); body.height = Inches(5.7)
        elif layout in ("below2", "below") and images:
            body.width = Inches(10.20); body.height = Inches(2.5 if layout == "below2" else 2.9)
        else:
            body.width = Inches(10.20); body.height = Inches(5.7)
        tf = body.text_frame; tf.word_wrap = True
        for i, b in enumerate(bullets):
            lvl, txt = b if isinstance(b, tuple) else (0, b)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = lvl
            r = p.add_run(); r.text = txt; _ea(r)            # ■/●/→ はレイアウトの箇条書きに従う
        images = images or []; caps = caps or []
        if placeholder is not None:
            self._ph(s, 5.55, 1.95, 4.95, 3.7, placeholder)
        elif layout == "right" and images:
            self._pic(s, images[0], 5.55, 1.95, 4.95, 4.4)
            if caps:
                self._cap(s, 5.55, 6.45, 4.95, caps[0])
        elif layout == "below2" and images:
            cw = (SW - 1.0) / 2
            for i, im in enumerate(images[:2]):
                self._pic(s, im, 0.45 + i * (cw + 0.1), 3.75, cw, 2.6)
                if i < len(caps):
                    self._cap(s, 0.45 + i * (cw + 0.1), 6.5, cw, caps[i])
        elif images:                                          # below（1枚を下に大きく）
            self._pic(s, images[0], 1.6, 4.05, 7.6, 2.85)
            if caps:
                self._cap(s, 1.6, 7.05, 7.6, caps[0])
        return s

    def slide(self, chapter, subtitle, keymsg, bullets, images=None, caps=None, placeholder=None):
        """新フォーマット（卒論1-1スライド準拠）:
          ヘッダ=2段（章=小14pt ＋ 小見出し=大, タイトルPHに）／キーメッセージ=全幅バナー／
          本文=左、図=右（画像があれば左右2分割）。subtitle=None なら章名のみ1行（02/05/06 等）。
        """
        LN = RGBColor(0xEA, 0xEF, 0xF7); N2 = RGBColor(0x1B, 0x2A, 0x5A); KR = RGBColor(0xB0, 0x30, 0x10)
        s = self.prs.slides.add_slide(self._layout())
        # --- ヘッダ（タイトルプレースホルダに2行）---
        title = s.shapes.title; tf = title.text_frame
        title.left = Inches(0.25); title.top = Inches(0.26); title.width = Inches(10.36); title.height = Inches(1.06)
        for r in list(tf.paragraphs[0].runs):
            r._element.getparent().remove(r._element)
        for ep in tf.paragraphs[1:]:
            ep._p.getparent().remove(ep._p)
        p0 = tf.paragraphs[0]; p0.space_before = Pt(0); p0.space_after = Pt(0)
        if subtitle:
            r = p0.add_run(); r.text = chapter; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = NAVY; _ea(r)
            p1 = tf.add_paragraph(); p1.space_before = Pt(0); p1.space_after = Pt(0)
            r = p1.add_run(); r.text = subtitle; r.font.bold = True; r.font.color.rgb = NAVY; _ea(r)  # サイズはPH既定（大）を継承
        else:
            r = p0.add_run(); r.text = chapter; r.font.bold = True; r.font.color.rgb = NAVY; _ea(r)
        # --- キーメッセージ（全幅バナー）---
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30), Inches(1.46), Inches(10.23), Inches(0.80))
        box.fill.solid(); box.fill.fore_color.rgb = LN; box.line.fill.background(); box.shadow.inherit = False
        try: box.adjustments[0] = 0.08
        except Exception: pass
        bt = box.text_frame; bt.word_wrap = True; bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bt.margin_left = Pt(10); bt.margin_right = Pt(10); bt.margin_top = Pt(2); bt.margin_bottom = Pt(2)
        bp = bt.paragraphs[0]
        r = bp.add_run(); r.text = "結論  "; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = KR; _ea(r)
        r = bp.add_run(); r.text = keymsg; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = N2; _ea(r)
        # --- 本文（左）＋ 図 / プレースホルダ（右）---
        has_r = bool(images) or (placeholder is not None)
        body = s.placeholders[1]
        body.left = Inches(0.30); body.top = Inches(2.42)
        body.width = Inches(4.95 if has_r else 10.20); body.height = Inches(4.7)
        tb = body.text_frame; tb.word_wrap = True
        for i, b in enumerate(bullets):
            lvl, txt = b if isinstance(b, tuple) else (0, b)
            p = tb.paragraphs[0] if i == 0 else tb.add_paragraph()
            p.level = lvl
            r = p.add_run(); r.text = txt; _ea(r)
        images = images or []; caps = caps or []
        if placeholder is not None:
            self._ph(s, 5.55, 2.50, 4.95, 3.9, placeholder)
        elif len(images) == 1:
            self._pic(s, images[0], 5.55, 2.50, 4.95, 4.0)
            if caps:
                self._cap(s, 5.55, 6.6, 4.95, caps[0])
        elif len(images) >= 2:
            for i, im in enumerate(images[:2]):
                self._pic(s, im, 5.55, 2.46 + i * 2.16, 4.95, 1.98)
                if i < len(caps):
                    self._cap(s, 5.55, 2.46 + i * 2.16 + 1.99, 4.95, caps[i])
        return s

    def divider(self, title):
        """中表紙（章区切り）。レイアウト '中表紙' を使う。"""
        s = self.prs.slides.add_slide(self._layout("中表紙"))
        if s.shapes.title is not None:
            s.shapes.title.text = title
        return s

    # ---------- 仕上げ ----------
    def save(self, out_path):
        """元テンプレの本文スライド(2..)を削除してから保存。
        ※新スライドは末尾に足してあるので、削除後の並び＝表紙・Index・新本文。"""
        sldIdLst = self.prs.slides._sldIdLst
        ids = list(sldIdLst)
        for i in sorted(range(self._keep, self._n0), reverse=True):
            rId = ids[i].get(qn("r:id"))
            self.prs.part.drop_rel(rId)
            sldIdLst.remove(ids[i])
        self.prs.save(out_path)
        return out_path
